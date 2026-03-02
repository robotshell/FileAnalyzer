import pdfplumber
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import zipfile, logging

def extract_text(file_path):
    text, metadata = "", []

    try:
        # PDF
        if file_path.endswith(".pdf"):
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ""
                if pdf.metadata:
                    metadata.extend([f"{k}: {v}" for k, v in pdf.metadata.items()])
                # Extraer comentarios ocultos
                if hasattr(pdf, "annots") and pdf.annots:
                    for annot in pdf.annots:
                        text += f"\n[PDF Comment]: {annot.get('contents','')}\n"

        # DOCX / DOCM
        elif file_path.endswith((".docx", ".docm")):
            doc = Document(file_path)
            for p in doc.paragraphs:
                text += p.text + "\n"
            # Metadatos
            if doc.core_properties.author:
                metadata.append(f"Author: {doc.core_properties.author}")
            if hasattr(doc, 'part') and hasattr(doc.part, 'vba_project'):
                metadata.append("Macros detected!")

        # XLSX / XLSM
        elif file_path.endswith((".xlsx", ".xlsm")):
            wb = load_workbook(file_path, data_only=True)
            for sheet in wb:
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell:
                            text += str(cell) + " "
            if file_path.endswith(".xlsm"):
                metadata.append("Macros detected!")

        # PPTX / PPTM
        elif file_path.endswith((".pptx", ".pptm")):
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            if file_path.endswith(".pptm"):
                metadata.append("Macros detected!")

    except Exception as e:
        logging.warning(f"Failed to extract {file_path}: {e}")

    return text, metadata
