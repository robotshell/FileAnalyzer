#!/usr/bin/env python3

import requests, re, json, sys, os, argparse, tempfile
from tqdm import tqdm
from colorama import Fore, Style, init

import pdfplumber
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

init(autoreset=True)

# =========================
# REGEX RULES
# =========================

EMAIL_REGEX = r"[a-zA-Z0-9._%+-]+@company\.com"
IBAN_REGEX = r"\b[A-Z]{2}\d{2}[A-Z0-9]{11,30}\b"
AWS_KEY = r"AKIA[0-9A-Z]{16}"
JWT = r"eyJ[A-Za-z0-9_-]+\.[A-Za-z0-9_-]+\.[A-Za-z0-9_-]+"

REGEX_RULES = {
    "Email detected": EMAIL_REGEX,
    "IBAN detected": IBAN_REGEX,
    "AWS Key detected": AWS_KEY,
    "JWT detected": JWT
}

# =========================
# RISK
# =========================

def calculate_risk(score):
    if score >= 60:
        return "HIGH"
    elif score >= 30:
        return "MEDIUM"
    return "LOW"

# =========================
# DOWNLOAD
# =========================

def download_file(url, timeout, max_size):
    try:
        r = requests.get(url, timeout=timeout, stream=True)
        if r.status_code != 200:
            return None

        size = int(r.headers.get("Content-Length", 0))
        if size > max_size:
            return None

        suffix = os.path.splitext(url)[1]
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)

        for chunk in r.iter_content(1024):
            tmp.write(chunk)
        tmp.close()
        return tmp.name

    except Exception:
        return None

# =========================
# TEXT EXTRACTION
# =========================

def extract_text(file_path):
    text, metadata = "", []

    try:
        if file_path.endswith(".pdf"):
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ""
                if pdf.metadata:
                    metadata.extend([f"{k}: {v}" for k, v in pdf.metadata.items()])

        elif file_path.endswith(".docx"):
            doc = Document(file_path)
            for p in doc.paragraphs:
                text += p.text + "\n"
            if doc.core_properties.author:
                metadata.append(f"Author metadata: {doc.core_properties.author}")

        elif file_path.endswith(".xlsx"):
            wb = load_workbook(file_path, data_only=True)
            for sheet in wb:
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell:
                            text += str(cell) + " "

        elif file_path.endswith(".pptx"):
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

    except Exception:
        pass

    return text, metadata

# =========================
# ANALYSIS
# =========================

def analyze(text, keywords):
    findings, score = [], 0

    for kw in keywords:
        if kw.lower() in text.lower():
            findings.append(f"Keyword found: {kw}")
            score += 10

    for name, regex in REGEX_RULES.items():
        matches = set(re.findall(regex, text))
        for m in matches:
            findings.append(f"{name}: {m}")
            score += 25

    return findings, score

# =========================
# MAIN
# =========================

def main(args):
    with open(args.urls) as f:
        urls = [u.strip() for u in f if u.strip()]

    with open(args.keywords) as f:
        keywords = [k.strip() for k in f if k.strip()]

    results = []

    print(Fore.CYAN + f"\n[+] Processing {len(urls)} files...\n")

    for url in tqdm(urls, desc="Scanning", ncols=80):
        tmp = download_file(url, args.timeout, args.max_size)
        if not tmp:
            continue

        text, metadata = extract_text(tmp)
        os.unlink(tmp)

        findings, score = analyze(text, keywords)

        for m in metadata:
            findings.append(m)
            score += 5

        if not findings:
            continue

        risk = calculate_risk(score)

        if args.silent and risk != "HIGH":
            continue

        result = {
            "url": url,
            "risk": risk,
            "score": score,
            "findings": findings
        }
        results.append(result)

        color = Fore.RED if risk == "HIGH" else Fore.YELLOW if risk == "MEDIUM" else Fore.GREEN
        print(color + f"\n[{risk}] {url} (score: {score})")
        for f in findings:
            print("  └─", f)

        if args.poc:
            os.makedirs("poc", exist_ok=True)
            fname = os.path.join("poc", re.sub(r'\W+', '_', url) + ".txt")
            with open(fname, "w") as p:
                p.write(url + "\n\n")
                p.write("\n".join(findings))
                p.write("\n\n--- TEXT EXTRACT ---\n")
                p.write(text[:1000])

    if args.json:
        print(json.dumps(results, indent=2))

# =========================
# ARGPARSE
# =========================

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Sensitive File Exposure Scanner")
    parser.add_argument("urls", help="File with URLs")
    parser.add_argument("keywords", help="Keywords file")
    parser.add_argument("--silent", action="store_true", help="Show only HIGH risk")
    parser.add_argument("--json", action="store_true", help="JSON output")
    parser.add_argument("--poc", action="store_true", help="Generate PoC files")
    parser.add_argument("--timeout", type=int, default=10)
    parser.add_argument("--max-size", type=int, default=10_000_000)

    args = parser.parse_args()
    main(args)
